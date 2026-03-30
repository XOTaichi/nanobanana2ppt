from __future__ import annotations

import math
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

from background_remover import BriaRMBG2Remover
from path_utils import load_manifest, resolve_local_path
from ppt_style_utils import estimate_font_size, infer_alignment, pick_text_color, pixels_to_inches


def preferred_visual_path(element: dict, manifest_dir: Path, prepared_dir: Path) -> Path | None:
    element_id = element.get("id")
    if element_id:
        tight_path = prepared_dir / f"{element_id}_tight.png"
        if tight_path.exists():
            return tight_path
        raw_rmbg_path = prepared_dir / "rmbg_raw" / f"{element_id}_nobg.png"
        if raw_rmbg_path.exists():
            return raw_rmbg_path
    return resolve_local_path(manifest_dir, element.get("nobg_path") or element.get("crop_path"))


def add_picture(slide, image_path: Path, bbox: list[int]) -> None:
    slide.shapes.add_picture(str(image_path), Inches(pixels_to_inches(bbox[0])), Inches(pixels_to_inches(bbox[1])), Inches(pixels_to_inches(max(2, bbox[2] - bbox[0]))), Inches(pixels_to_inches(max(2, bbox[3] - bbox[1]))))


def add_textbox(slide, element: dict, manifest_dir: Path) -> None:
    text = (element.get("text") or "").strip()
    if not text:
        return
    bbox = list(element["bbox"])
    textbox = slide.shapes.add_textbox(Inches(pixels_to_inches(bbox[0])), Inches(pixels_to_inches(bbox[1])), Inches(pixels_to_inches(max(4, bbox[2] - bbox[0]))), Inches(pixels_to_inches(max(4, bbox[3] - bbox[1]))))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = infer_alignment(text, bbox)
    run = paragraph.runs[0]
    run.font.size = Pt(estimate_font_size(bbox))
    run.font.bold = len(text) <= 18 or bbox[3] - bbox[1] >= 36
    crop_path = resolve_local_path(manifest_dir, element.get("crop_path"))
    run.font.color.rgb = pick_text_color(crop_path) if crop_path else RGBColor(60, 45, 30)
    textbox.line.fill.background()
    textbox.fill.background()


def add_white_background(slide, width_px: int, height_px: int) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(pixels_to_inches(width_px)), Inches(pixels_to_inches(height_px)))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()


def add_title(slide, title: str, width_px: int) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.2), Inches(max(2.0, pixels_to_inches(width_px) - 0.7)), Inches(0.4))
    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.runs[0].font.size = Pt(20)
    paragraph.runs[0].font.bold = True
    paragraph.runs[0].font.color.rgb = RGBColor(40, 40, 40)
    title_box.line.fill.background()
    title_box.fill.background()


def add_contained_picture(slide, image_path: Path, left_in: float, top_in: float, width_in: float, height_in: float) -> None:
    with Image.open(image_path) as img:
        scale = min(width_in / pixels_to_inches(img.width), height_in / pixels_to_inches(img.height))
        render_width = pixels_to_inches(img.width) * scale
        render_height = pixels_to_inches(img.height) * scale
    add_picture(slide, image_path, [int((left_in + (width_in - render_width) / 2) * 96), int((top_in + (height_in - render_height) / 2) * 96), int((left_in + (width_in + render_width) / 2) * 96), int((top_in + (height_in + render_height) / 2) * 96)])


def prepare_visuals(visual_elements: list[dict], manifest_dir: Path, prepared_dir: Path, use_rmbg: bool, rmbg_model_path: Path | None, min_alpha_pixels: int, min_alpha_ratio: float) -> list[dict]:
    remover = BriaRMBG2Remover(rmbg_model_path, prepared_dir / "rmbg_raw") if use_rmbg else None
    prepared: list[dict] = []
    for element in visual_elements:
        source_path = preferred_visual_path(element, manifest_dir, prepared_dir)
        if source_path is None:
            continue
        if remover is None:
            prepared.append({**element, "prepared_visual_path": str(source_path)})
            continue
        rgba = Image.open(remover.remove_background(Image.open(source_path), element["id"])).convert("RGBA")
        alpha_bin = rgba.getchannel("A").point(lambda v: 255 if v >= 24 else 0)
        alpha_bbox = alpha_bin.getbbox()
        if alpha_bbox is None:
            continue
        non_transparent = sum(1 for value in alpha_bin.getdata() if value > 0)
        bbox_area = max(1, (alpha_bbox[2] - alpha_bbox[0]) * (alpha_bbox[3] - alpha_bbox[1]))
        if non_transparent < min_alpha_pixels or non_transparent / bbox_area < min_alpha_ratio:
            continue
        tight_path = prepared_dir / f"{element['id']}_tight.png"
        prepared_dir.mkdir(parents=True, exist_ok=True)
        rgba.crop(alpha_bbox).save(tight_path)
        bbox = list(element["bbox"])
        prepared.append({**element, "bbox": [bbox[0] + alpha_bbox[0], bbox[1] + alpha_bbox[1], bbox[0] + alpha_bbox[2], bbox[1] + alpha_bbox[3]], "prepared_visual_path": str(tight_path)})
    return prepared


def add_crop_gallery(prs: Presentation, visual_elements: list[dict], manifest_dir: Path, width_px: int, height_px: int) -> None:
    crops = [(element.get("id") or "crop", resolve_local_path(manifest_dir, element.get("crop_path"))) for element in visual_elements]
    crops = [(label, path) for label, path in crops if path and path.exists()]
    if not crops:
        return
    cols, rows, page_size = 4, 3, 12
    for page_idx in range(math.ceil(len(crops) / page_size)):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_white_background(slide, width_px, height_px)
        add_title(slide, f"SAM3 Visual Crops {page_idx + 1}/{math.ceil(len(crops) / page_size)}", width_px)
        for item_idx, (label, image_path) in enumerate(crops[page_idx * page_size : (page_idx + 1) * page_size]):
            row, col = divmod(item_idx, cols)
            cell_left, cell_top = 0.3 + col * 2.4, 0.55 + row * 2.8
            frame = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(cell_left), Inches(cell_top), Inches(2.2), Inches(2.55))
            frame.fill.solid()
            frame.fill.fore_color.rgb = RGBColor(248, 248, 248)
            frame.line.color.rgb = RGBColor(220, 220, 220)
            add_contained_picture(slide, image_path, cell_left + 0.08, cell_top + 0.06, 2.04, 2.0)
            caption = slide.shapes.add_textbox(Inches(cell_left + 0.05), Inches(cell_top + 2.12), Inches(2.1), Inches(0.18))
            paragraph = caption.text_frame.paragraphs[0]
            paragraph.text = label
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.runs[0].font.size = Pt(9)
            paragraph.runs[0].font.color.rgb = RGBColor(70, 70, 70)
            caption.line.fill.background()
            caption.fill.background()


def export_ppt(manifest_path: Path, output_pptx: Path, background_mode: str, text_only: bool, use_rmbg: bool, rmbg_model_path: Path | None, min_alpha_pixels: int, min_alpha_ratio: float) -> Path:
    manifest = load_manifest(manifest_path)
    manifest_dir = manifest_path.parent
    input_image = resolve_local_path(manifest_dir, manifest.get("input_image"))
    restored = resolve_local_path(manifest_dir, manifest.get("artifacts", {}).get("restored_background"))
    if input_image is None or not input_image.exists():
        raise RuntimeError("input_image not found in manifest or on disk.")
    width_px = manifest["image_size"]["width"]
    height_px = manifest["image_size"]["height"]
    prs = Presentation()
    prs.slide_width = Inches(pixels_to_inches(width_px))
    prs.slide_height = Inches(pixels_to_inches(height_px))
    original_slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_white_background(original_slide, width_px, height_px)
    add_title(original_slide, "Original Image", width_px)
    add_contained_picture(original_slide, input_image, 0.35, 0.75, max(1.0, pixels_to_inches(width_px) - 0.7), max(1.0, pixels_to_inches(height_px) - 1.1))
    editable_slide = prs.slides.add_slide(prs.slide_layouts[6])
    if background_mode == "restored" and restored and restored.exists():
        add_picture(editable_slide, restored, [0, 0, width_px, height_px])
    else:
        add_white_background(editable_slide, width_px, height_px)
    elements = manifest.get("elements", [])
    text_elements = [element for element in elements if element.get("element_type") == "text"]
    visual_elements = [element for element in elements if element.get("element_type") == "visual"]
    for element in text_elements:
        add_textbox(editable_slide, element, manifest_dir)
    if not text_only:
        for element in prepare_visuals(visual_elements, manifest_dir, manifest_dir / "prepared_visuals", use_rmbg, rmbg_model_path, min_alpha_pixels, min_alpha_ratio):
            add_picture(editable_slide, Path(element["prepared_visual_path"]), list(element["bbox"]))
        add_crop_gallery(prs, visual_elements, manifest_dir, width_px, height_px)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))
    return output_pptx
