#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import os
import sys
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Inches, Pt

from artifact_builder import save_crops
from baidu_ocr_client import detect_text_elements
from env_utils import load_dotenv
from path_utils import manifest_default_dir
from ppt_style_utils import infer_alignment, pick_text_color, pixels_to_inches


BASE_DIR = Path(__file__).resolve().parent


def add_white_text_mask(slide, bbox: list[int], padding_px: int = 1) -> None:
    x0, y0, x1, y1 = max(0, bbox[0] - padding_px), max(0, bbox[1] - padding_px), max(bbox[0] + 2, bbox[2] + padding_px), max(bbox[1] + 2, bbox[3] + padding_px)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(pixels_to_inches(x0)), Inches(pixels_to_inches(y0)), Inches(pixels_to_inches(max(2, x1 - x0))), Inches(pixels_to_inches(max(2, y1 - y0))))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()
    shape.shadow.inherit = False


def add_editable_text(slide, element: dict, manifest_dir: Path, font_name: str | None, font_size: float | None) -> None:
    bbox = list(element["bbox"])
    textbox = slide.shapes.add_textbox(Inches(pixels_to_inches(bbox[0])), Inches(pixels_to_inches(bbox[1])), Inches(pixels_to_inches(max(4, bbox[2] - bbox[0]))), Inches(pixels_to_inches(max(4, bbox[3] - bbox[1]))))
    frame = textbox.text_frame
    frame.word_wrap = True
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    paragraph = frame.paragraphs[0]
    paragraph.text = (element.get("text") or "").strip()
    paragraph.alignment = infer_alignment(paragraph.text, bbox)
    run = paragraph.runs[0]
    if font_name:
        run.font.name = font_name
    run.font.size = Pt(font_size if font_size is not None else max(8.0, min(32.0, (bbox[3] - bbox[1]) * 0.42)))
    run.font.bold = len(paragraph.text) <= 18 or bbox[3] - bbox[1] >= 36
    crop_path = Path(element["crop_path"]) if element.get("crop_path") else None
    if crop_path and not crop_path.is_absolute():
        crop_path = (manifest_dir / crop_path).resolve()
    run.font.color.rgb = pick_text_color(crop_path) if crop_path and crop_path.exists() else RGBColor(60, 45, 30)
    textbox.line.fill.background()
    textbox.fill.background()


def build_text_only_ppt(input_image: Path, output_pptx: Path, text_elements: list[dict], font_name: str | None, font_size: float | None) -> None:
    with Image.open(input_image) as img:
        width_px, height_px = img.size
    prs = Presentation()
    prs.slide_width = Inches(pixels_to_inches(width_px))
    prs.slide_height = Inches(pixels_to_inches(height_px))
    original = prs.slides.add_slide(prs.slide_layouts[6])
    original.shapes.add_picture(str(input_image), 0, 0, Inches(pixels_to_inches(width_px)), Inches(pixels_to_inches(height_px)))
    editable = prs.slides.add_slide(prs.slide_layouts[6])
    editable.shapes.add_picture(str(input_image), 0, 0, Inches(pixels_to_inches(width_px)), Inches(pixels_to_inches(height_px)))
    for element in text_elements:
        add_white_text_mask(editable, list(element["bbox"]))
    for element in text_elements:
        add_editable_text(editable, element, output_pptx.parent, font_name, font_size)
    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Run the OCR-only text replacement pipeline.")
    parser.add_argument("input_png", type=Path, help="Input PNG path")
    parser.add_argument("--output-dir", type=Path, default=BASE_DIR / "output_text_only", help="Base output directory")
    parser.add_argument("--output", type=Path, default=None, help="Output PPTX path")
    parser.add_argument("--dotenv", type=Path, default=BASE_DIR / ".env", help="Path to .env file")
    parser.add_argument("--ocr-api-key", default=None)
    parser.add_argument("--font-name", default=None)
    parser.add_argument("--font-size", type=float, default=None)
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    load_dotenv(args.dotenv)
    input_image = args.input_png.resolve()
    if not input_image.exists():
        print(f"Input image not found: {input_image}", file=sys.stderr)
        return 2
    ocr_api_key = args.ocr_api_key or os.environ.get("BAIDU_API_KEY")
    if not ocr_api_key:
        print("Missing BAIDU_API_KEY / --ocr-api-key", file=sys.stderr)
        return 2
    output_dir = manifest_default_dir(input_image, args.output_dir.resolve())
    output_dir.mkdir(parents=True, exist_ok=True)
    output_pptx = args.output.resolve() if args.output else (output_dir / "text_only_editable.pptx").resolve()
    image = Image.open(input_image).convert("RGB")
    text_records = detect_text_elements(input_image, ocr_api_key)
    save_crops(image, text_records, output_dir / "crops", remove_bg=False, rmbg_model_path=None)
    text_elements = [record.to_dict() for record in text_records]
    build_text_only_ppt(input_image, output_pptx, text_elements, args.font_name, args.font_size)
    (output_dir / "text_manifest.json").write_text(json.dumps({"input_image": str(input_image), "image_size": {"width": image.width, "height": image.height}, "elements": text_elements, "counts": {"text": len(text_elements), "total": len(text_elements)}, "artifacts": {"pptx": str(output_pptx)}}, ensure_ascii=False, indent=2), encoding="utf-8")
    print(output_pptx)
    return 0


if __name__ == "__main__":
    raise SystemExit(main())

